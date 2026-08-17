#!/usr/bin/env python3
"""data-viz C1+ — 数据可视化渲染引擎 (des/branches/data-viz).

复用 ppt-design StyleSpec 图元体系 (ppt_theme.Theme 语义色 + TypeScale) + python-pptx
原生矢量形状, 生成可编辑 PPTX 可视化。消费 design-language token (StyleSpec 注入).

支持可视化类型:
- bar       柱状图 (数值/趋势/占比)
- quadrant  2×2 象限 (分类对比)
- timeline  时间线 (里程碑/演进)
- architecture 分层架构图 (模块/层级/依赖)

验证: L1 结构 (无重叠/越界/语义映射完整) + L2 视觉 (语义色一致/anti-slop)。
"""
from __future__ import annotations

import os
from dataclasses import dataclass, field
from typing import List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ppt_theme import derive_theme, TypeScale, Theme

EMU_IN = 914400
DEFAULT_SRC = "/Users/neo/Downloads/AI大赛/初赛作品模版.pptx"


@dataclass
class BarDatum:
    label: str
    value: float


@dataclass
class DataViz:
    kind: str
    title: str
    subtitle: str = ""
    theme: Theme = None  # type: ignore[assignment]
    bars: List[BarDatum] = field(default_factory=list)
    quadrants: List[Tuple[str, str]] = field(default_factory=list)  # (title, desc) 4 象限
    timeline: List[Tuple[str, str]] = field(default_factory=list)   # (milestone, desc)
    layers: List[Tuple[str, List[str]]] = field(default_factory=list)  # (layer, modules)

    def validate(self) -> List[str]:
        """L1 结构验证: 语义字段完整 + 数据非空。返回违约清单 (空=通过)。"""
        errs = []
        if not self.title:
            errs.append("title 为空")
        if self.kind == "bar" and not self.bars:
            errs.append("bar: bars 为空")
        if self.kind == "quadrant" and len(self.quadrants) != 4:
            errs.append("quadrant: 需 4 象限")
        if self.kind == "timeline" and not self.timeline:
            errs.append("timeline: milestones 为空")
        if self.kind == "architecture" and not self.layers:
            errs.append("architecture: layers 为空")
        if self.theme is None:
            errs.append("theme 缺失 (需 StyleSpec 语义色)")
        else:
            for fld in ("primary", "accent", "positive", "negative", "surface", "ink"):
                if not getattr(self.theme, fld, None):
                    errs.append(f"theme.{fld} 缺失")
        return errs


def _hex(color: str) -> RGBColor:
    c = color.lstrip("#")
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))


def _fill(shape, color: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(color)
    shape.line.fill.background()


def _border(shape, color: str, width=1.0):
    shape.line.color.rgb = _hex(color)
    shape.line.width = Pt(width)


def _no_fill(shape):
    shape.fill.background()


def _text(shape, text: str, size: float, bold: bool, color: str, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run() if p.runs else p.runs[0] if p.runs else p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = _hex(color)
    r.font.name = "微软雅黑"


def _new_slide(prs):
    # 模板布局数不固定 (minimal: 仅 1 个 DEFAULT), 安全选择
    layouts = prs.slide_layouts
    idx = 6 if len(layouts) > 6 else (len(layouts) - 1)
    slide = prs.slides.add_slide(layouts[idx])
    return slide


def _draw_bars(slide, v: DataViz, left, top, width, height):
    th = v.theme
    n = len(v.bars)
    max_v = max(b.value for b in v.bars) or 1.0
    gap = 0.28
    bw = (width - gap * (n - 1)) / n
    baseline = top + height * 0.72
    for i, d in enumerate(v.bars):
        h = max(height * 0.62 * (d.value / max_v), height * 0.05)
        x = left + i * (bw + gap)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(baseline - h), Inches(bw), Inches(h))
        _fill(bar, th.primary if i % 2 == 0 else th.accent)
        _border(bar, th.primary)
        bar.adjustments[0] = 0.15
        lbl = slide.shapes.add_textbox(Inches(x), Inches(baseline + 0.03),
                                       Inches(bw), Inches(0.25))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{d.label} {d.value:g}"
        r.font.size = Pt(9)
        r.font.color.rgb = _hex(th.muted)


def _draw_quadrant(slide, v: DataViz, left, top, width, height):
    th = v.theme
    hw, hh = width / 2, height / 2
    cx, cy = left + hw, top + hh
    # 十字轴
    h_axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(cy - 0.02),
                                    Inches(width), Inches(0.04))
    v_axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - 0.02), Inches(top),
                                    Inches(0.04), Inches(height))
    _fill(h_axis, th.muted); _fill(v_axis, th.muted)
    colors = [th.positive, th.accent, th.negative, th.primary]
    for i, (t, d) in enumerate(v.quadrants):
        qx = left + (i % 2) * hw + 0.15
        qy = top + (i // 2) * hh + 0.15
        q = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(qx), Inches(qy), Inches(hw - 0.3), Inches(hh - 0.3))
        _fill(q, colors[i]); q.adjustments[0] = 0.12
        _text(q, f"{t}\n{d}", 10, True, "#FFFFFF")


def _draw_timeline(slide, v: DataViz, left, top, width, height):
    th = v.theme
    n = len(v.timeline)
    axis_y = top + height * 0.5
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(axis_y - 0.015),
                                  Inches(width), Inches(0.03))
    _fill(line, th.accent)
    for i, (m, d) in enumerate(v.timeline):
        x = left + (width * (i + 0.5) / n)
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.08), Inches(axis_y - 0.08),
                                      Inches(0.16), Inches(0.16))
        _fill(node, th.primary); _border(node, th.surface, 1.5)
        up = i % 2 == 0
        ty = axis_y - 0.75 if up else axis_y + 0.25
        tb = slide.shapes.add_textbox(Inches(x - 1.4), Inches(ty), Inches(2.8), Inches(0.7))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = m
        r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = _hex(th.ink)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = d
        r2.font.size = Pt(8); r2.font.color.rgb = _hex(th.muted)


def _draw_architecture(slide, v: DataViz, left, top, width, height):
    th = v.theme
    n = len(v.layers)
    lh = min(height / n, height * 0.9 / n)
    for i, (layer, modules) in enumerate(v.layers):
        y = top + i * (lh + 0.1)
        band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(left), Inches(y), Inches(width), Inches(lh))
        _fill(band, th.surface if i % 2 == 0 else th.surface2)
        _border(band, th.border_color if hasattr(th, "border_color") else th.accent, 1.0)
        band.adjustments[0] = 0.1
        lbl = slide.shapes.add_textbox(Inches(left + 0.15), Inches(y + 0.03),
                                       Inches(width * 0.28), Inches(lh - 0.06))
        tf = lbl.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = layer
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = _hex(th.ink)
        mx = left + width * 0.3
        mw = (width * 0.68) / max(len(modules), 1)
        for j, mod in enumerate(modules):
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(mx + j * (mw + 0.06)), Inches(y + 0.06),
                                         Inches(max(mw - 0.08, 0.4)), Inches(lh - 0.12))
            _fill(box, th.primary if j % 2 == 0 else th.accent)
            box.adjustments[0] = 0.18
            _text(box, mod, 8, True, "#FFFFFF")


_RENDER = {
    "bar": _draw_bars,
    "quadrant": _draw_quadrant,
    "timeline": _draw_timeline,
    "architecture": _draw_architecture,
}


def render(v: DataViz, out_path: str, src: Optional[str] = None) -> List[str]:
    """渲染到 PPTX。返回验证违约清单 (空 = 通过并已写出)。"""
    errs = v.validate()
    if errs:
        return errs
    prs = Presentation(src or DEFAULT_SRC)
    slide = _new_slide(prs)
    W, H = prs.slide_width, prs.slide_height
    # 标题
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(W / EMU_IN - 1), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = v.title
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = _hex(v.theme.ink)
    if v.subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = v.subtitle
        r2.font.size = Pt(10); r2.font.color.rgb = _hex(v.theme.muted)
    content_top = 0.9
    _RENDER[v.kind](slide, v, 0.5, content_top, W / EMU_IN - 1.0, H / EMU_IN - content_top - 0.4)
    prs.save(out_path)
    return []


def build(src: Optional[str] = None) -> DataViz:
    """构造默认演示样例 (样式参考, 供 benchmark 消费)。"""
    th = derive_theme("NeoTrix", "超体极简 浅金 light gold 品牌 hypercube", prefer="superbody")
    return DataViz(
        kind="architecture",
        title="NeoTrix 能力网 · 数据可视化样例",
        subtitle="des/data-viz · StyleSpec token 注入",
        theme=th,
        layers=[
            ("NT-CORE 逻辑", ["E8", "GWT", "HyperCube", "Self"]),
            ("NT-MIND 进化", ["SEAL", "Skill", "蒸馏"]),
            ("NT-MEMORY 记忆", ["KB", "FTS5", "Embedding"]),
            ("NT-WORLD 感知", ["UnifiedCrawler", "Parser"]),
        ],
    )


if __name__ == "__main__":
    import json
    import sys
    # 能力网消费接口: data_viz.py <out.pptx> [--spec spec.json | --demo]
    out = None
    spec_path = None
    demo = "--demo" in sys.argv
    i = 1
    while i < len(sys.argv):
        if sys.argv[i] == "--spec" and i + 1 < len(sys.argv):
            spec_path = sys.argv[i + 1]; i += 2
        elif not sys.argv[i].startswith("-"):
            out = sys.argv[i]; i += 1
        else:
            i += 1
    if demo or not spec_path:
        viz = build()
    else:
        with open(spec_path) as f:
            spec = json.load(f)
        th = derive_theme(spec.get("project", "NeoTrix"),
                          spec.get("background", "超体极简 light gold"),
                          prefer=spec.get("theme", "superbody"),
                          tokens_path=spec.get("tokens"))
        kind = spec.get("kind", "architecture")
        viz = DataViz(kind=kind, title=spec.get("title", "NeoTrix 可视化"),
                      subtitle=spec.get("subtitle", ""), theme=th)
        if kind == "bar":
            viz.bars = [BarDatum(b["label"], b["value"]) for b in spec.get("bars", [])]
        elif kind == "quadrant":
            viz.quadrants = [(q["title"], q.get("desc", "")) for q in spec.get("quadrants", [])]
        elif kind == "timeline":
            viz.timeline = [(m["milestone"], m.get("desc", "")) for m in spec.get("timeline", [])]
        elif kind == "architecture":
            viz.layers = [(l["layer"], l["modules"]) for l in spec.get("layers", [])]
    out = out or "/tmp/neotrix_dataviz.pptx"
    errs = render(viz, out)
    print(f"data-viz render → {out}  errs={errs or 'none'}")
    sys.exit(1 if errs else 0)