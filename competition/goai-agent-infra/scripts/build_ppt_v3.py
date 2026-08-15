#!/usr/bin/env python3
"""初赛方案 PPT v3 — 设计大师排版。

基于 ppt_theme 引擎：从项目背景推导主题（意识流：navy 底 + amber 沉淀/自进化 +
primary blue 执行 + positive green 价值 + negative red 痛点），并按专业字体层级
（kicker 小节标签 / lead 导语 / h1 卡片标题 / body 正文 / support 辅助）排版。

语义化配色（跨页一致）：
- primary blue   = 执行 / 协作 / 方案主体
- accent amber   = 沉淀 / 自进化 / 知识结晶 / 高光
- positive green = 价值 / 通过 / 开源
- negative red   = 痛点 / 风险 / 打回
- support teal/cyan = 辅助 / 次级
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from ppt_theme import derive_theme, t, TypeScale

SRC = "/Users/neo/Downloads/AI大赛/初赛作品模版.pptx"
OUT = "/Users/neo/Downloads/neotrix/competition/goai-agent-infra/初赛方案_NeoTrix研发闭环.pptx"

PROJECT = "NeoTrix 研发闭环"
BACKGROUND = ("AI-native 软件研发全流程多 Agent 协同系统：GWT 注意力调度编排 5 职能 Agent，"
              "E8 确定性推理定位根因，SelfTest T1-T3 质量门禁，experience-tree 复盘自动结晶 Skill，"
              "MIT 开源、4240 项测试背书")
THEME = derive_theme(PROJECT, BACKGROUND)

FONT = "微软雅黑"
_ts = TypeScale()


def _set_ea(run, name=FONT):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)


def _set_track(run, em=0.0, size=12.0):
    """letter-spacing: em → spc (1/100 pt)。"""
    if not em:
        return
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(round(em * size * 100))))


def add_shape(slide, shp, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.10):
    sp = slide.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(line_w)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def put_text(sp, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             margin=0.06, tracking_em=0.0):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            txt, sz, bd, cl = line[0], line[1], line[2], line[3]
            lh = line[4] if len(line) > 4 else _ts.line_dense
        else:
            txt, sz, bd, cl, lh = line, _ts.body, False, THEME.ink, _ts.line_body
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = lh
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = bd
        run.font.color.rgb = RGBColor.from_string(cl)
        run.font.name = FONT
        _set_ea(run)
        _set_track(run, tracking_em, sz)
    return sp


def box(slide, x, y, w, h, lines, fill=THEME.surface, line=None, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, tracking_em=0.0):
    sp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)
    return put_text(sp, lines, align, anchor, tracking_em=tracking_em)


def arrow(slide, x, y, w, h, dir=MSO_SHAPE.RIGHT_ARROW, fill=None):
    return add_shape(slide, dir, x, y, w, h, fill or THEME.primary)


def pill(slide, x, y, w, h, text, fill):
    """小节标签胶囊（kicker，带字距）。"""
    return box(slide, x, y, w, h, [t(THEME, "kicker", text)], fill, tracking_em=0.12)


def support_bar(slide, x, y, w, h, lines, fill=None):
    box(slide, x, y, w, h, lines, fill or THEME.surface2, align=PP_ALIGN.LEFT)


def card(slide, x, y, w, title, items, strip):
    """标准卡片：色带头 + 条目区。"""
    box(slide, x, y, w, 0.5, [t(THEME, "h1", title)], strip, tracking_em=0.08)
    box(slide, x, y + 0.57, w, 1.95, items, THEME.surface, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)


def clear_content(slide, keep_top_in=1.25):
    threshold = int(keep_top_in * 914400)
    for shp in list(slide.shapes):
        if shp.top is not None and shp.top >= threshold:
            shp._element.getparent().remove(shp._element)


# ================= 内容页 =================

def draw_s5(s):
    th = THEME
    # 导语条
    box(s, 0.55, 1.35, 12.2, 0.42,
        [t(th, "lead", "目标用户：企业研发团队 / 开源维护者    ·   真实场景：GitHub Issue → 聚合 → 定位 → 修复 → 审查 → 沉淀")],
        th.base, tracking_em=0.02)
    # 三栏
    pill(s, 0.55, 1.92, 3.9, 0.42, "现状痛点", th.negative)
    pains = [
        [t(th, "body-bold", "① 缺陷多源分散"), t(th, "micro", "Issue/日志/反馈/CI 分散，无统一证据集")],
        [t(th, "body-bold", "② 根因依赖人工"), t(th, "micro", "资深工程师经验，不可复现")],
        [t(th, "body-bold", "③ 质量无法量化"), t(th, "micro", "修复无门禁，复盘不沉淀")],
    ]
    for i, pl in enumerate(pains):
        y = 2.44 + i * 1.13
        box(s, 0.55, y, 3.9, 1.0, pl, th.surface, align=PP_ALIGN.LEFT)
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.55, y, 0.07, 1.0, th.negative)
    arrow(s, 4.6, 2.9, 0.7, 0.5, MSO_SHAPE.RIGHT_ARROW)
    # 中：方案
    pill(s, 5.4, 1.92, 2.6, 0.42, "NeoTrix 方案", th.primary)
    mid = [
        t(th, "h2", "AgentTeams"),
        t(th, "white-body", "Manager + A1-A5"),
        t(th, "white-body", "5 职能 Agent"),
        t(th, "white-body", "确定性推理"),
        t(th, "accent", "自进化沉淀"),
    ]
    box(s, 5.4, 2.44, 2.6, 2.39, mid, th.base, line=th.primary)
    add_shape(s, MSO_SHAPE.RECTANGLE, 5.4, 2.44, 2.6, 0.07, th.accent)
    arrow(s, 8.15, 2.9, 0.7, 0.5, MSO_SHAPE.RIGHT_ARROW, th.positive)
    # 右：价值
    pill(s, 8.95, 1.92, 3.8, 0.42, "可量化价值", th.positive)
    vals = [
        [t(th, "body-bold", "定位 30-60min → 5-10min", color=th.ink), t(th, "micro", "效率提升指标")],
        [t(th, "body-bold", "SelfTest T1-T3 门禁", color=th.ink), t(th, "micro", "修复成功率 / 门禁通过率")],
        [t(th, "body-bold", "复盘自动结晶 Skill", color=th.ink), t(th, "micro", "知识沉淀速率")],
    ]
    for i, pl in enumerate(vals):
        y = 2.44 + i * 1.13
        box(s, 8.95, y, 3.8, 1.0, pl, th.surface, align=PP_ALIGN.LEFT)
        add_shape(s, MSO_SHAPE.RECTANGLE, 8.95, y, 0.07, 1.0, th.positive)
    # 辅助条
    support_bar(s, 0.55, 5.95, 12.2, 0.98,
        [t(th, "support", "行业可复制性：任何有 Issue+CI+代码仓库的组织可复制，可迁移 IT 服务/嵌入式/金融科技研发线"),
         t(th, "support", "差异化（对照 OpenHands/SWE-agent）：确定性推理可复现 + T1-T3 生产门禁 + experience-tree 自进化沉淀")])


def draw_s7(s):
    th = THEME
    layers = [
        ("任务输入层", "GitHub Issue / 日志 / 用户反馈 / CI 告警", th.base),
        ("AgentTeams 编排层", "Manager（拆解/委派/追踪） · A1-A5 Worker · Matrix 房间全程可见 · MinIO 共享上下文", th.primary),
        ("Skill 能力层", "官方用云 Skills（resourcecenter-search / ecs-diagnose / network-reachability / sas-overview）+ dev-implementer / rev-officer / experience-tree / mcp-gateway", th.support1),
        ("MCP / 工具层", "GitHub · CI · 监控 · 云产品   （Higress 统一网关托管凭证，Worker 不持真实 key）", th.support2),
        ("证据与治理层", "MinIO 共享状态 · trace_data · SHA-256 审计链 · KB 向量+BM25 检索", th.accent),
    ]
    y = 1.46
    for i, (tt, d, c) in enumerate(layers):
        box(s, 0.55, y, 3.0, 0.92, [t(th, "h2", tt)], c, tracking_em=0.1)
        box(s, 3.7, y, 9.05, 0.92, [t(th, "body", d)],
            th.surface if i % 2 else th.surface2, align=PP_ALIGN.LEFT)
        if i < len(layers) - 1:
            arrow(s, 12.62, y + 0.97, 0.3, 0.14, MSO_SHAPE.DOWN_ARROW)
        y += 1.05
    support_bar(s, 0.55, 6.65, 12.2, 0.45,
        [t(th, "support", "选型必要性：AgentTeams=必选协同基点 · Higress=统一网关+凭证隔离 · MinIO=共享上下文降 Token 消耗")])


def draw_s9(s):
    th = THEME
    box(s, 0.55, 1.46, 2.3, 0.8,
        [t(th, "h2", "Human"), t(th, "white-body", "Matrix 房间全程可见可干预")], th.negative)
    arrow(s, 2.95, 1.71, 0.5, 0.4, MSO_SHAPE.RIGHT_ARROW)
    box(s, 3.55, 1.46, 2.6, 0.8,
        [t(th, "h2", "Manager 主控"), t(th, "white-body", "拆解·委派·追踪"), t(th, "white-body", "不自主批准高风险动作")],
        th.base, line=th.accent)
    box(s, 6.85, 1.46, 5.9, 0.8,
        [t(th, "h2", "上下文传递载体", color=th.ink), t(th, "body", "MinIO shared/tasks/<id>/ · KB 结构化中间结论 · Matrix 时间线"),
         t(th, "micro", "Worker 无状态可替换")], th.surface, align=PP_ALIGN.LEFT)
    # A1-A4 流水线
    agents = [
        ("A1 采集者", "聚合/去重/分级", "NT-WORLD"),
        ("A2 诊断者", "E8 根因定位", "NT-MIND"),
        ("A3 实施者", "TDD 修复", "NT-ACT"),
        ("A4 审查者", "独立验证·门禁", "NT-SHIELD"),
    ]
    x = 0.55
    for i, (n, d, dom) in enumerate(agents):
        box(s, x, 2.70, 2.85, 1.15,
            [t(th, "h2", n), t(th, "white-body", d), t(th, "micro", dom, color="DCE6FF")], th.primary)
        if i < 3:
            arrow(s, x + 2.85, 3.10, 0.35, 0.4, MSO_SHAPE.RIGHT_ARROW)
        x += 3.2
    box(s, 7.4, 4.05, 5.4, 0.5, [t(th, "negative", "A4 审查不通过 → 打回 A3 重试（最多 N 轮，封顶升级人工）")],
        th.surface, align=PP_ALIGN.LEFT)
    # A5 + KB
    box(s, 0.55, 4.77, 2.85, 1.15,
        [t(th, "h2", "A5 沉淀者"), t(th, "white-body", "experience-tree 吸收"), t(th, "white-body", "KB 经验 + Skill 回灌")], th.accent)
    arrow(s, 3.45, 5.07, 0.4, 0.4, MSO_SHAPE.RIGHT_ARROW)
    box(s, 3.9, 4.77, 8.85, 1.15,
        [t(th, "h2", "知识库 KB", color=th.ink), t(th, "body", "experience 命名空间 · 向量+BM25 混合检索"),
         t(th, "accent", "复盘自动结晶为 Skill → 回灌 Manager 拆解逻辑（自进化闭环）")],
        th.surface, align=PP_ALIGN.LEFT)
    # 辅助条
    support_bar(s, 0.55, 6.15, 12.2, 1.0,
        [t(th, "support", "状态流转：Task 状态机（待办→拆解→执行→审查→沉淀）；审查不通过打回 A3，封顶升级人工"),
         t(th, "support", "安全边界：高风险动作（改生产/大重构/删数据）→ 人工审批 · Worker 仅持 Higress 消费者令牌 · SHA-256 审计链")])


def draw_s11(s):
    th = THEME
    groups = [
        ("官方用云 Skills（必选）", [
            ("alibabacloud-resourcecenter-search 资源搜索", True),
            ("alibabacloud-ecs-diagnose 实例诊断", True),
            ("alibabacloud-network-reachability-analysis 网络分析", True),
            ("alibabacloud-sas-overview 安全态势", True)], th.primary),
        ("编码 / 实施", [("dev-implementer  TDD 红-绿-重构", False), ("github-operations  Issue/分支/PR", False), ("repair-healer  自愈修复与恢复验证", False)], th.support1),
        ("审查 / 治理", [("rev-officer  全量审查 D1-D63+S1-S7", False), ("gov-steward  合规策略", False), ("mcp-gateway  MCP 工具聚合", False)], th.support2),
        ("知识 / 吸收", [("experience-tree  五阶段吸收协议", False), ("sg-diagnostician  元认知自审计", False), ("KB RAG  向量+BM25 混合检索", False)], th.accent),
    ]
    pos = [(0.55, 1.46), (6.85, 1.46), (0.55, 4.00), (6.85, 4.00)]
    for (tt, items, c), (x, y) in zip(groups, pos):
        card(s, x, y, 5.95, tt, [t(th, "body", it) for it, _ in items], c)
    support_bar(s, 0.55, 6.55, 12.2, 0.55,
        [t(th, "support", "生命周期：SKILL.md 装载 Worker 工作区 → Manager 按需分发 → 版本/发布/回滚（安全审核→灰度→审计） · 4 个官方 Skill 串联即排障链路")])


def draw_s13(s):
    th = THEME
    quads = [
        ("可运行性", ["cargo build/check/test 双验证", "Docker 部署 · install.sh", "CI workflows（.github）"], th.primary),
        ("运行证据", ["4240 项测试 · SelfTest T1-T3 三层接线", "日志 / Trace / Metrics 全记录", "trace_data + SHA-256 广播审计链"], th.support1),
        ("可观测", ["Skill/MCP/RAG/LLM 全链路 Trace", "Log 结构化关联 TraceId", "Metrics：修复成功率/时延/Token/Tool 成功率"], th.support2),
        ("安全治理", ["RAG：KB 向量+BM25，证据强制溯源", "权限矩阵 · 审批 · 回滚 · 审计", "gitleaks 密钥扫描 · 零 unsafe"], th.positive),
    ]
    pos = [(0.55, 1.46), (6.85, 1.46), (0.55, 4.00), (6.85, 4.00)]
    for (tt, items, c), (x, y) in zip(quads, pos):
        card(s, x, y, 5.95, tt, [t(th, "body", it) for it in items], c)
    support_bar(s, 0.55, 6.55, 12.2, 0.55,
        [t(th, "support", "云产品选型：Higress（网关）统一入口+凭证托管，可替换性/迁移成本已论证 · 全链路证据可审计、可回滚")])


def draw_s15(s):
    th = THEME
    steps = [
        ("可复用成果", ["Skill 体系独立发布", "mcp-gateway 网关", "KB 检索层 · E8 推理内核"], th.primary),
        ("接口契约与文档", ["README · 部署说明", "开源协议 · 示例配置", "测试方法 · 复现步骤"], th.support1),
        ("协议与依赖披露", ["MIT", "第三方依赖全披露", "商业 API / 闭源模型边界"], th.support2),
        ("社区共建", ["4240 测试背书", "AI 工程/研发效能社区", "长期维护与迭代"], th.positive),
    ]
    x = 0.55
    for i, (tt, d, c) in enumerate(steps):
        card(s, x, 1.55, 2.9, tt, [t(th, "body", it) for it in d], c)
        if i < 3:
            arrow(s, x + 2.9, 2.55, 0.3, 0.4, MSO_SHAPE.RIGHT_ARROW)
        x += 3.2
    support_bar(s, 0.55, 5.0, 12.2, 1.0,
        [t(th, "support", "开源范围：核心引擎 + Skill + 网关 + 示例与运行报告，均可复用可验证"),
         t(th, "support", "合规：数据来源与授权边界、第三方依赖、商业 API、闭源模型使用范围逐一披露")])


def draw_s17(s):
    th = THEME
    miles = [
        ("初赛 8.16", ["提交：简介 + 方案 PPT"], th.primary),
        ("复赛名单 8.24", ["Top30 入围"], th.support1),
        ("复赛 9.3", ["可执行 AgentTeams 代码包", "端到端 Demo / 视频"], th.support2),
        ("决赛名单 9.10", ["Top15 入围"], th.accent),
        ("决赛 9.22", ["现场路演 + Demo", "代码仓库最终版"], th.positive),
    ]
    x = 0.55
    for i, (tt, d, c) in enumerate(miles):
        card(s, x, 1.55, 2.2, tt, [t(th, "body", it) for it in d], c)
        if i < 4:
            arrow(s, x + 2.2, 1.72, 0.3, 0.35, MSO_SHAPE.RIGHT_ARROW)
        x += 2.5
    support_bar(s, 0.55, 4.45, 12.2, 0.9,
        [t(th, "support", "复赛工程化：AgentTeams 本地 install.sh → K8s helm · GitHub Issue 端到端 Demo · SWE-bench-style 评测 · 官方 Skills 接入")])
    box(s, 0.55, 5.5, 12.2, 1.15,
        [t(th, "body-bold", "评估指标：修复成功率 · 端到端时延 · Token 成本 · 门禁通过率 · 经验沉淀速率"),
         t(th, "body", "风险控制：Demo 环境不确定性（网络/凭证）→ Mock 与真实接入共用同一 Schema · 评审口径严格对齐评分维度")],
        th.surface, align=PP_ALIGN.LEFT)


DRAW = {5: draw_s5, 7: draw_s7, 9: draw_s9, 11: draw_s11, 13: draw_s13, 15: draw_s15, 17: draw_s17}

FILL = {
    1: {"文本框 18": "NeoTrix 研发闭环\n软件研发全流程多 Agent 协同系统\nAgent Infra · 方向三 · 初赛方案"},
    2: {
        "Text 7": "NeoTrix 研发闭环",
        "Text 12": "缺陷链路碎片化：Issue/日志/反馈多源分散，根因定位依赖资深工程师经验，修复质量无量化门禁，复盘经验不沉淀。",
        "Text 17": "AgentTeams Manager-Workers 编排 5 职能 Agent（采集→诊断→实施→审查→沉淀），E8 确定性推理 + experience-tree 自进化。",
        "Text 22": "① E8 确定性 64 态推理内核，根因定位可复现；② SelfTest T1-T3 生产门禁（4240 测试）+ 复盘自动结晶为 Skill。",
        "Text 27": "MIT 开源、Skill 体系独立可复用、RAG/可观测可迁移到任意多 Agent 生产系统。",
        "Text 32": "方案设计完成；7 域架构 23.6 万行 + 4240 测试 + MCP/Skill/KB 基础设施可支撑复赛 Demo。",
    },
    19: {"文本框 37": "团队名称：NeoTrix（个人参赛 · 郑州大学）\n成员：Asher（郑州大学·软件开发/AI Agent 方向）\n团队分工：Asher — 主控/架构 · Agent/Skill 工程 · Demo/验证（个人参赛，一人承担全流程）\n团队成果：NeoTrix 开源项目（4240 项测试、RQGM 论文 arXiv:2606.26294）；软件研发全流程多 Agent 协同系统\n作品链接：github.com/neo-trixs/NeoTrix"},
}


def set_text(shape, text):
    tf = shape.text_frame
    lines = text.split("\n")
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lines[0]
        for r in p0.runs[1:]:
            r.text = ""
    else:
        p0.add_run().text = lines[0]
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.add_run().text = line


def main():
    print(f"推导主题: {THEME.name}  命中信号: {THEME.signals}")
    prs = Presentation(SRC)
    for idx, shapes_map in FILL.items():
        slide = prs.slides[idx - 1]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name in shapes_map:
                set_text(shape, shapes_map[shape.name])
    for slide_idx, draw_fn in DRAW.items():
        s = prs.slides[slide_idx - 1]
        clear_content(s)
        draw_fn(s)
    prs.save(OUT)
    print("已生成:", OUT)


if __name__ == "__main__":
    sys.exit(main())