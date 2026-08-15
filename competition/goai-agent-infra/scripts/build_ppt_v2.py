#!/usr/bin/env python3
"""初赛方案 PPT v2 — 图为主、文字辅助、原生矢量流程图排版。

设计原则：
1. 删除模板占位图片（通用素材不传达信息），改为原生矢量图（圆角矩形+箭头）绘制真实方案内容
2. 每页 = 大图讲清楚 + 底部/侧边 1-2 行文字辅助，杜绝遮挡
3. 矢量图形在 OfficeCLI 渲染链路中任意缩放清晰
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SRC = "/Users/neo/Downloads/AI大赛/初赛作品模版.pptx"
OUT = "/Users/neo/Downloads/neotrix/competition/goai-agent-infra/初赛方案_NeoTrix研发闭环.pptx"

# ---- 配色（与模板 navy 主题统一）----
NAVY   = "1B1F3B"
BLUE   = "2E5BFF"
LBLUE  = "3B82F6"
TEAL   = "0EA5E9"
CYAN   = "06B6D4"
LIGHT  = "F1F5F9"
LGRAY  = "E2E8F0"
GRAY   = "64748B"
DARK   = "0F172A"
RED    = "EF4444"
GREEN  = "10B981"
ORANGE = "F59E0B"
WHITE  = "FFFFFF"

FONT = "微软雅黑"


def _set_ea(run, name=FONT):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)


def add_shape(slide, shp, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.12):
    sp = slide.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(line_w)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def put_text(sp, lines, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, margin=0.05, line_spacing=1.0):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            txt, sz, bd, cl = line
        else:
            txt, sz, bd, cl = line, size, bold, color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = bd
        run.font.color.rgb = RGBColor.from_string(cl)
        run.font.name = FONT
        _set_ea(run)
    return sp


def box(slide, x, y, w, h, lines, fill=BLUE, color=WHITE, size=11, bold=True,
        line=None, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    sp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)
    return put_text(sp, lines, size, color, bold, align, anchor)


def arrow(slide, x, y, w, h, dir=MSO_SHAPE.RIGHT_ARROW, fill=GRAY):
    return add_shape(slide, dir, x, y, w, h, fill)


def label(slide, x, y, w, h, text, fill=None, color=DARK, size=11, bold=True):
    return box(slide, x, y, w, h, [text], fill or LIGHT, color, size, bold)


def clear_content(slide, keep_top_in=1.25):
    """删除 top >= keep_top_in 的所有 shape（保留页眉），删除模板占位图片与文本框。"""
    threshold = int(keep_top_in * 914400)
    for shp in list(slide.shapes):
        if shp.top is not None and shp.top >= threshold:
            shp._element.getparent().remove(shp._element)


# ============ 各内容页绘图 ============

def draw_s5(s):
    """场景与价值：痛点 → 5 Agent 闭环 → 可量化价值"""
    label(s, 0.55, 1.35, 12.2, 0.42,
          "目标用户：企业研发团队 / 开源维护者    ·   真实场景：GitHub Issue → 聚合 → 定位 → 修复 → 审查 → 沉淀",
          NAVY, WHITE, 12)
    # 左：痛点
    label(s, 0.55, 1.86, 3.9, 0.4, "现状痛点", RED, WHITE, 12)
    pains = [
        (("① 缺陷多源分散", 12, True, DARK), ("Issue/日志/反馈/CI 分散，无统一证据集", 10, False, GRAY)),
        (("② 根因依赖人工", 12, True, DARK), ("资深工程师经验，不可复现", 10, False, GRAY)),
        (("③ 质量无法量化", 12, True, DARK), ("修复无门禁，复盘不沉淀", 10, False, GRAY)),
    ]
    for i, (t, d) in enumerate(pains):
        box(s, 0.55, 2.36 + i * 1.18, 3.9, 1.0, [t, d], LIGHT, DARK, align=PP_ALIGN.LEFT)
    arrow(s, 4.6, 2.9, 0.7, 0.5, MSO_SHAPE.RIGHT_ARROW, BLUE)
    # 中：方案
    label(s, 5.4, 1.86, 2.6, 0.4, "NeoTrix 方案", BLUE, WHITE, 12)
    box(s, 5.4, 2.36, 2.6, 2.4,
        [("AgentTeams", 13, True, WHITE), ("Manager + A1-A5", 10, False, WHITE),
         ("5 职能 Agent", 10, False, WHITE), ("确定性推理", 10, False, WHITE),
         ("自进化沉淀", 10, False, WHITE)], NAVY, WHITE, line=BLUE)
    arrow(s, 8.15, 2.9, 0.7, 0.5, MSO_SHAPE.RIGHT_ARROW, GREEN)
    # 右：价值
    label(s, 8.95, 1.86, 3.8, 0.4, "可量化价值", GREEN, WHITE, 12)
    vals = [
        (("定位耗时 30-60min → 5-10min", 12, True, DARK), ("效率提升指标", 10, False, GRAY)),
        (("SelfTest T1-T3 质量门禁", 12, True, DARK), ("修复成功率/门禁通过率", 10, False, GRAY)),
        (("复盘自动结晶为 Skill", 12, True, DARK), ("知识沉淀速率", 10, False, GRAY)),
    ]
    for i, (t, d) in enumerate(vals):
        box(s, 8.95, 2.36 + i * 1.18, 3.8, 1.0, [t, d], LIGHT, DARK, align=PP_ALIGN.LEFT)
    # 底部辅助文字
    box(s, 0.55, 5.95, 12.2, 0.95,
        [("行业可复制性：任何有 Issue+CI+代码仓库的组织可复制，可迁移 IT 服务/嵌入式/金融科技研发线", 11, False, DARK),
         ("差异化（对照 OpenHands/SWE-agent）：确定性推理可复现 + T1-T3 生产门禁 + experience-tree 自进化沉淀", 11, False, DARK)],
        LGRAY, DARK, align=PP_ALIGN.LEFT)


def draw_s7(s):
    """方案总览：四层架构图"""
    layers = [
        ("任务输入层", "GitHub Issue / 日志 / 用户反馈 / CI 告警", NAVY),
        ("AgentTeams 编排层", "Manager（拆解/委派/追踪） · A1-A5 Worker · Matrix 房间全程可见 · MinIO 共享上下文", BLUE),
        ("Skill 能力层", "官方用云 Skills（resourcecenter-search / ecs-diagnose / network-reachability / sas-overview）+ dev-implementer / rev-officer / experience-tree / mcp-gateway", LBLUE),
        ("MCP / 工具层", "GitHub · CI · 监控 · 云产品   （Higress 统一网关托管凭证，Worker 不持真实 key）", TEAL),
        ("证据与治理层", "MinIO 共享状态 · trace_data · SHA-256 审计链 · KB 向量+BM25 检索", CYAN),
    ]
    y = 1.46
    for i, (t, d, c) in enumerate(layers):
        box(s, 0.55, y, 3.0, 0.92, [t], c, WHITE, 12)
        box(s, 3.7, y, 9.05, 0.92, [d], LIGHT if i % 2 else LGRAY, DARK, 11, False, align=PP_ALIGN.LEFT)
        if i < len(layers) - 1:
            arrow(s, 12.62, y + 0.97, 0.3, 0.14, MSO_SHAPE.DOWN_ARROW, BLUE)
        y += 1.05
    box(s, 0.55, 6.65, 12.2, 0.45,
        [("选型必要性：AgentTeams=必选协同基点 · Higress=统一网关+凭证隔离 · MinIO=共享上下文降 Token 消耗", 10, False, DARK)],
        LGRAY, DARK)


def draw_s9(s):
    """多 Agent 协同：闭环流程图"""
    box(s, 0.55, 1.46, 2.3, 0.8, [("Human", 12, True, WHITE), ("Matrix 房间全程可见可干预", 9, False, WHITE)], RED, WHITE)
    arrow(s, 2.95, 1.71, 0.5, 0.4, MSO_SHAPE.RIGHT_ARROW, BLUE)
    box(s, 3.55, 1.46, 2.6, 0.8, [("Manager 主控", 12, True, WHITE), ("拆解·委派·追踪", 9, False, WHITE), ("不自主批准高风险动作", 9, False, WHITE)], NAVY, WHITE)
    arrow(s, 6.25, 1.71, 0.5, 0.4, MSO_SHAPE.RIGHT_ARROW, BLUE)
    box(s, 6.85, 1.46, 5.9, 0.8,
        [("上下文传递载体", 12, True, DARK), ("MinIO shared/tasks/<id>/ · KB 结构化中间结论 · Matrix 时间线", 9.5, False, GRAY),
         ("Worker 无状态可替换", 9.5, False, GRAY)], LIGHT, DARK)
    # A1-A4 横向流水线
    agents = [
        ("A1 采集者", "聚合/去重/分级", "NT-WORLD"),
        ("A2 诊断者", "E8 根因定位", "NT-MIND"),
        ("A3 实施者", "TDD 修复", "NT-ACT"),
        ("A4 审查者", "独立验证·门禁", "NT-SHIELD"),
    ]
    x = 0.55
    for i, (n, d, dom) in enumerate(agents):
        box(s, x, 2.70, 2.85, 1.15, [n, d, dom], BLUE, WHITE, 12)
        if i < 3:
            arrow(s, x + 2.85, 3.10, 0.35, 0.4, MSO_SHAPE.RIGHT_ARROW, GRAY)
        x += 3.2
    # 审查打回反馈（文字说明 + 左箭头位于 A4 下方）
    box(s, 7.4, 4.05, 5.4, 0.5, [("A4 审查不通过 → 打回 A3 重试（最多 N 轮，封顶升级人工）", 10.5, True, RED)],
        LIGHT, RED, align=PP_ALIGN.LEFT)
    # A5 + KB 沉淀
    box(s, 0.55, 4.77, 2.85, 1.15, [("A5 沉淀者", 12, True, WHITE), ("experience-tree 吸收", 9.5, False, WHITE), ("KB 经验 + Skill 回灌", 9.5, False, WHITE)], ORANGE, WHITE)
    arrow(s, 3.45, 5.07, 0.4, 0.4, MSO_SHAPE.RIGHT_ARROW, GRAY)
    box(s, 3.9, 4.77, 8.85, 1.15,
        [("知识库 KB", 12, True, DARK), ("experience 命名空间 · 向量+BM25 混合检索", 9.5, False, GRAY),
         ("复盘自动结晶为 Skill → 回灌 Manager 拆解逻辑（自进化闭环）", 10, False, BLUE)], LIGHT, DARK, align=PP_ALIGN.LEFT)
    # 底部辅助
    box(s, 0.55, 6.15, 12.2, 1.0,
        [("状态流转：Task 状态机（待办→拆解→执行→审查→沉淀）；审查不通过打回 A3，封顶升级人工", 11, False, DARK),
         ("安全边界：高风险动作（改生产/大重构/删数据）→ 人工审批 · Worker 仅持 Higress 消费者令牌 · SHA-256 审计链", 11, False, DARK)],
        LGRAY, DARK, align=PP_ALIGN.LEFT)


def draw_s11(s):
    """Skill 体系：四类分组 2x2"""
    groups = [
        ("官方用云 Skills（必选）", ["alibabacloud-resourcecenter-search 资源搜索", "alibabacloud-ecs-diagnose 实例诊断", "alibabacloud-network-reachability-analysis 网络分析", "alibabacloud-sas-overview 安全态势"], BLUE),
        ("编码 / 实施", ["dev-implementer  TDD 红-绿-重构", "github-operations  Issue/分支/PR", "repair-healer  自愈修复与恢复验证"], LBLUE),
        ("审查 / 治理", ["rev-officer  全量审查 D1-D63+S1-S7", "gov-steward  合规策略", "mcp-gateway  MCP 工具聚合"], TEAL),
        ("知识 / 吸收", ["experience-tree  五阶段吸收协议", "sg-diagnostician  元认知自审计", "KB RAG  向量+BM25 混合检索"], CYAN),
    ]
    pos = [(0.55, 1.46), (6.85, 1.46), (0.55, 4.00), (6.85, 4.00)]
    for (t, items, c), (x, y) in zip(groups, pos):
        box(s, x, y, 5.95, 0.55, [t], c, WHITE, 12)
        box(s, x, y + 0.62, 5.95, 1.80, items, LIGHT, DARK, 10.5, False, anchor=MSO_ANCHOR.TOP)
    box(s, 0.55, 6.55, 12.2, 0.55,
        [("生命周期：SKILL.md 装载 Worker 工作区 → Manager 按需分发 → 版本/发布/回滚（安全审核→灰度→审计）", 11, False, DARK),
         ("复用价值：Skill 为任务能力抽象层，跨 Agent 跨场景复用；4 个官方 Skill 串联即排障链路", 11, False, DARK)],
        LGRAY, DARK, align=PP_ALIGN.LEFT)


def draw_s13(s):
    """工程落地：四象限"""
    quads = [
        ("可运行性", ["cargo build/check/test 双验证", "Docker 部署 · install.sh", "CI workflows（.github）"], LBLUE),
        ("运行证据", ["4240 项测试 · SelfTest T1-T3 三层接线", "日志 / Trace / Metrics 全记录", "trace_data + SHA-256 广播审计链"], TEAL),
        ("可观测", ["Skill/MCP/RAG/LLM 全链路 Trace", "Log 结构化关联 TraceId", "Metrics：修复成功率/时延/Token/Tool 成功率"], CYAN),
        ("安全治理", ["RAG：KB 向量+BM25，证据强制溯源", "权限矩阵 · 审批 · 回滚 · 审计", "gitleaks 密钥扫描 · 零 unsafe"], GREEN),
    ]
    pos = [(0.55, 1.46), (6.85, 1.46), (0.55, 4.00), (6.85, 4.00)]
    for (t, items, c), (x, y) in zip(quads, pos):
        box(s, x, y, 5.95, 0.55, [t], c, WHITE, 12)
        box(s, x, y + 0.62, 5.95, 1.80, items, LIGHT, DARK, 10.5, False, anchor=MSO_ANCHOR.TOP)
    box(s, 0.55, 6.55, 12.2, 0.55,
        [("云产品选型：Higress（网关）统一入口+凭证托管，可替换性/迁移成本已论证 · 全链路证据可审计、可回滚", 11, False, DARK)],
        LGRAY, DARK)


def draw_s15(s):
    """开源：四步横向流程"""
    steps = [
        ("可复用成果", ["Skill 体系独立发布", "mcp-gateway 网关", "KB 检索层 · E8 推理内核"], BLUE),
        ("接口契约与文档", ["README · 部署说明", "开源协议 · 示例配置", "测试方法 · 复现步骤"], LBLUE),
        ("协议与依赖披露", ["MIT", "第三方依赖全披露", "商业 API / 闭源模型边界"], TEAL),
        ("社区共建", ["4240 测试背书", "AI 工程/研发效能社区", "长期维护与迭代"], CYAN),
    ]
    x = 0.55
    for i, (t, d, c) in enumerate(steps):
        box(s, x, 1.55, 2.9, 0.55, [t], c, WHITE, 12)
        box(s, x, 2.17, 2.9, 2.2, d, LIGHT, DARK, 10.5, False, anchor=MSO_ANCHOR.TOP)
        if i < 3:
            arrow(s, x + 2.9, 2.55, 0.3, 0.4, MSO_SHAPE.RIGHT_ARROW, GRAY)
        x += 3.2
    box(s, 0.55, 5.0, 12.2, 1.0,
        [("开源范围：核心引擎 + Skill + 网关 + 示例与运行报告，均可复用可验证", 11, False, DARK),
         ("合规：数据来源与授权边界、第三方依赖、商业 API、闭源模型使用范围逐一披露", 11, False, DARK)],
        LGRAY, DARK, align=PP_ALIGN.LEFT)


def draw_s17(s):
    """落地计划：里程碑时间线"""
    miles = [
        ("初赛 8.16", ["提交：简介 + 方案 PPT"], BLUE),
        ("复赛名单 8.24", ["Top30 入围"], LBLUE),
        ("复赛 9.3", ["可执行 AgentTeams 代码包", "端到端 Demo / 视频"], TEAL),
        ("决赛名单 9.10", ["Top15 入围"], ORANGE),
        ("决赛 9.22", ["现场路演 + Demo", "代码仓库最终版"], GREEN),
    ]
    x = 0.55
    for i, (t, d, c) in enumerate(miles):
        box(s, x, 1.55, 2.2, 0.55, [t], c, WHITE, 12)
        box(s, x, 2.17, 2.2, 1.9, d, LIGHT, DARK, 10.5, False, anchor=MSO_ANCHOR.TOP)
        if i < 4:
            arrow(s, x + 2.2, 1.72, 0.3, 0.35, MSO_SHAPE.RIGHT_ARROW, GRAY)
        x += 2.5
    box(s, 0.55, 4.45, 12.2, 0.9,
        [("复赛工程化：AgentTeams 本地 install.sh → K8s helm · GitHub Issue 端到端 Demo · SWE-bench-style 评测 · 官方 Skills 接入", 11, False, DARK)],
        LGRAY, DARK)
    box(s, 0.55, 5.5, 12.2, 1.15,
        [("评估指标：修复成功率 · 端到端时延 · Token 成本 · 门禁通过率 · 经验沉淀速率", 11, False, DARK),
         ("风险控制：Demo 环境不确定性（网络/凭证）→ Mock 与真实接入共用同一 Schema · 评审口径严格对齐评分维度", 11, False, DARK)],
        LIGHT, DARK, align=PP_ALIGN.LEFT)


DRAW = {5: draw_s5, 7: draw_s7, 9: draw_s9, 11: draw_s11, 13: draw_s13, 15: draw_s15, 17: draw_s17}

# ---- 封面 / P0 / 团队页文本（与 v1 一致）----
FILL = {
    1: {"文本框 18": "NeoTrix 研发闭环\n软件研发全流程多 Agent 协同系统\nAgent Infra · 方向三 · 初赛方案"},
    2: {
        "Text 7": "NeoTrix 研发闭环",
        "Text 12": "缺陷链路碎片化：Issue/日志/反馈多源分散，根因定位依赖资深工程师经验，修复质量无量化门禁，复盘经验不沉淀。",
        "Text 17": "AgentTeams Manager-Workers 编排 5 职能 Agent（采集→诊断→实施→审查→沉淀），E8 确定性推理 + experience-tree 自进化。",
        "Text 22": "① E8 确定性 64 态推理内核，根因定位可复现；② SelfTest T1-T3 生产门禁（4240 测试）+ 复盘自动结晶为 Skill。",
        "Text 27": "MIT 开源、Skill 体系独立可复用、RAG/可观测可迁移到任意多 Agent 生产系统。",
        "Text 32": "方案设计完成；7 域架构 23.6 万行 + 4240 测试 + MCP/Skill/KB 基础设施可支撑复赛 Demo。",
    },
    19: {"文本框 37": "团队名称：NeoTrix（个人参赛 · 郑州大学）\n成员：Asher（郑州大学·软件开发/AI Agent 方向）\n团队分工：Asher — 主控/架构 · Agent/Skill 工程 · Demo/验证（个人参赛，一人承担全流程）\n团队成果：NeoTrix 开源项目（4240 项测试、RQGM 论文 arXiv:2606.26294）；软件研发全流程多 Agent 协同系统\n作品链接：github.com/neo-trixs/NeoTrix"},
}


def set_text(shape, text):
    tf = shape.text_frame
    lines = text.split("\n")
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lines[0]
        for r in p0.runs[1:]:
            r.text = ""
    else:
        p0.add_run().text = lines[0]
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.add_run().text = line


def main():
    prs = Presentation(SRC)
    for idx, shapes_map in FILL.items():
        slide = prs.slides[idx - 1]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name in shapes_map:
                set_text(shape, shapes_map[shape.name])
    for slide_idx, draw_fn in DRAW.items():
        s = prs.slides[slide_idx - 1]
        clear_content(s)
        draw_fn(s)
    prs.save(OUT)
    print("已生成:", OUT)


if __name__ == "__main__":
    sys.exit(main())
