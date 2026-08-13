#!/usr/bin/env python3
"""用初赛方案真实内容填充官方 PPT 模板（保留模板版式，只改文本）。"""
import sys
from pptx import Presentation

SRC = "/Users/neo/Downloads/AI大赛/初赛作品模版.pptx"
OUT = "/Users/neo/Downloads/neotrix/competition/goai-agent-infra/初赛方案_NeoTrix研发闭环.pptx"

# slide index (1-based) -> {shape_name: new_text}
CONTENT = {
    1: {
        "文本框 18": "NeoTrix 研发闭环\n软件研发全流程多 Agent 协同系统\nAgent Infra · 方向三 · 初赛方案",
    },
    2: {
        "Text 7": "NeoTrix 研发闭环",
        "Text 12": "缺陷链路碎片化：Issue/日志/反馈多源分散，根因定位依赖资深工程师经验，修复质量无量化门禁，复盘经验不沉淀。",
        "Text 17": "AgentTeams Manager-Workers 编排 5 职能 Agent（采集→诊断→实施→审查→沉淀），E8 确定性推理 + experience-tree 自进化。",
        "Text 22": "① E8 确定性 64 态推理内核，根因定位可复现；② SelfTest T1-T3 生产门禁（4240 测试）+ 复盘自动结晶为 Skill。",
        "Text 27": "MIT 开源、Skill 体系独立可复用、RAG/可观测可迁移到任意多 Agent 生产系统。",
        "Text 32": "方案设计完成；7 域架构 23.6 万行 + 4240 测试 + MCP/Skill/KB 基础设施可支撑复赛 Demo。",
    },
    5: {
        "文本框 4": "目标用户：企业研发团队 / 开源维护者。核心痛点：缺陷定位靠人工经验、修复质量无法量化验证、复盘经验散落。\n真实场景：GitHub Issue → 聚合 → E8 根因定位 → TDD 修复 → 独立审查 → 复盘沉淀。\n价值收益：定位耗时对标人工 30-60min 降至 5-10min、修复成功率、门禁通过率、知识沉淀条数。\n行业可复制性：任何有 Issue+CI+代码仓库的组织可复制；可迁移 IT 服务/嵌入式/金融科技研发线。",
        "文本框 6": "差异化对照 OpenHands/SWE-agent：确定性推理可复现 + T1-T3 生产门禁 + experience-tree 自进化沉淀。",
    },
    7: {
        "文本框 37": "架构：AgentTeams 编排层（Manager+A1-A5）→ Skill 能力层（官方用云 Skills + dev-implementer/rev-officer/experience-tree）→ MCP 工具层（GitHub/CI/监控，Higress 托管凭证）→ 证据治理层（MinIO 共享状态 + trace + SHA-256 审计链 + KB）。\n选型必要性：AgentTeams=必选协同基点；Higress=统一网关+凭证；MinIO=共享上下文降 Token。",
    },
    9: {
        "文本框 37": "Agent 分工：Manager(拆解/委派/追踪) + A1 采集 + A2 诊断 + A3 实施 + A4 审查 + A5 沉淀（见附录A）。\n任务拆解：collect→diagnose→implement→audit→distill 映射 AgentTeams task-management Skill。\n上下文传递：MinIO shared/tasks 共享工作区 + Matrix 房间时间线 + KB 结构化中间结论；Worker 无状态可替换。\n状态流转：Task 状态机 + 审查不通过打回 A3（重试封顶升级人工）。\n高风险动作：改生产/大重构/删数据→人工审批；Worker 不持真实凭证。",
    },
    11: {
        "文本框 37": "官方用云 Skills（必选）：alibabacloud-resourcecenter-search / ecs-diagnose / network-reachability-analysis / sas-overview / data-agent-skill，4 个串联即排障链路。\n核心 Skill：dev-implementer / rev-officer / repair-healer / experience-tree / mcp-gateway / github-operations（输入输出/依赖/失败处理/安全边界见附录B）。\n复用性：Skill 为任务能力抽象层，SKILL.md 装载 Worker 工作区，Manager 按需分发；版本/发布/回滚经安全审核→灰度→审计。",
    },
    13: {
        "文本框 37": "可运行性：cargo build/check/test 双验证、Docker 部署、CI workflows。\n运行证据：4240 测试、SelfTest T1-T3 三层接线、日志/Trace/Metrics 全记录。\n可观测：Skill/MCP/RAG/LLM 全链路 Trace + Log(TraceId 关联) + Metrics(修复成功率/时延/Token/Tool成功率)。\nRAG：KB nodes/edges + 向量 + BM25 混合检索，证据强制溯源。\n安全：权限矩阵、审批、回滚、审计、gitleaks、零 unsafe。",
    },
    15: {
        "文本框 37": "可复用成果：Skill 体系独立发布、mcp-gateway 网关、KB 检索层、E8 推理内核。\n接口契约与文档示例：README、部署说明、开源协议、示例配置、测试方法。\n协议与依赖：MIT；披露全部第三方依赖、商业 API 调用、闭源模型、数据授权边界。",
    },
    17: {
        "文本框 37": "当前进展：方案设计完成；7 域架构 23.6 万行 + 4240 测试 + MCP/Skill/KB 基础设施。\n里程碑：8.16 初赛提交→8.24 复赛名单→9.3 复赛提交→9.10 决赛名单→9.22 决赛答辩。\n复赛计划：AgentTeams 本地 install.sh→K8s helm；GitHub Issue 端到端 Demo；SWE-bench-style 评测；官方 Skills 接入。\n风险控制：Demo 环境不确定性→Mock+真实共用同一 Schema；评审口径→严格对齐评分维度。",
    },
    19: {
        "文本框 37": "成员背景：[姓名]（学校/公司·岗位·技能）。\n团队分工（≤3 人）：[姓名]—主控/架构；[姓名]—Agent/Skill 工程；[姓名]—Demo/验证。\n团队成果：NeoTrix 开源（4240 测试、RQGM 论文 arXiv:2606.26294）、过往获奖/项目。\n作品链接：github.com/neo-trixs/NeoTrix",
    },
}


def set_text(shape, text):
    tf = shape.text_frame
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if p.runs:
            p.runs[0].text = line
            for r in p.runs[1:]:
                r.text = ""
        else:
            p.add_run().text = line
    # remove leftover extra paragraphs
    while len(tf.paragraphs) > len(lines):
        el = tf.paragraphs[-1]._p
        el.getparent().remove(el)


def main():
    prs = Presentation(SRC)
    filled = []
    for idx, shapes_map in CONTENT.items():
        slide = prs.slides[idx - 1]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.name in shapes_map:
                set_text(shape, shapes_map[shape.name])
                filled.append(f"slide{idx}:{shape.name}")
    prs.save(OUT)
    print("已生成:", OUT)
    print(f"填充占位符 {len(filled)} 处")
    for f in filled:
        print(" ", f)


if __name__ == "__main__":
    sys.exit(main())
